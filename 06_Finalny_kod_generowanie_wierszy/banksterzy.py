import numpy as np
import pandas as pd
import matplotlib.pyplot as plt
import matplotlib.ticker as mtick

import io
import os
import re
import nltk
import scipy
import spacy
import ffmpeg
import pickle
import shutil
import tarfile
import pl_spacy_model
import speech_recognition 
import scipy.io.wavfile

from io import BytesIO
from fuzzywuzzy import fuzz
from base64 import b64decode
from google.colab import output
from stop_words import get_stop_words
from IPython.display import Javascript
from google.colab.output import eval_js
from scipy.io.wavfile import read as wav_read

from nltk.corpus import stopwords
from nltk.stem import WordNetLemmatizer
from sklearn.feature_selection import chi2
from sklearn.model_selection import train_test_split
from sklearn.feature_extraction.text import TfidfVectorizer
from sklearn.feature_extraction.text import TfidfTransformer, CountVectorizer

from itertools import chain
from pptx.util import Inches
from pptx.util import Cm, Pt
from pptx import Presentation  
from google.colab import files
from pd2ppt import df_to_table
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

RECORD = """
const sleep  = time => new Promise(resolve => setTimeout(resolve, time))
const b2text = blob => new Promise(resolve => {
  const reader = new FileReader() 
  reader.onloadend = e => resolve(e.srcElement.result)
  reader.readAsDataURL(blob)
})
var record = time => new Promise(async resolve => {
  stream = await navigator.mediaDevices.getUserMedia({ audio: true })
  recorder = new MediaRecorder(stream)
  chunks = []
  recorder.ondataavailable = e => chunks.push(e.data)
  recorder.start()
  await sleep(time)
  recorder.onstop = async ()=>{
    blob = new Blob(chunks)
    text = await b2text(blob)
    resolve(text)
  }
  recorder.stop()
})
"""

def record(sec=3):
  print("Mów teraz")
  display(Javascript(RECORD))
  sec += 1
  s = output.eval_js('record(%d)' % (sec*1000))
  print("Nagrywanie zakończone, dziękuję!")
  b = b64decode(s.split(',')[1])

  process = (ffmpeg
    .input('pipe:0')
    .output('pipe:1', format='wav')
    .run_async(pipe_stdin=True, pipe_stdout=True, pipe_stderr=True, quiet=True, overwrite_output=True)
  )
  output_fin, err = process.communicate(input=b)
  
  riff_chunk_size = len(output_fin) - 8
  # Break up the chunk size into four bytes, held in b.
  q = riff_chunk_size
  b_null = []
  for i in range(4):
      q, r = divmod(q, 256)
      b_null.append(r)

  # Replace bytes 4:8 in proc.stdout with the actual size of the RIFF chunk.
  riff = output_fin[:4] + bytes(b_null) + output_fin[8:]

  sr, audio = wav_read(io.BytesIO(riff))

  return audio, sr

def podyktuj_wers(sec):
  audio, sr = record(sec)
  scipy.io.wavfile.write('recording.wav', sr, audio)
  r = speech_recognition.Recognizer()
  filename = 'recording.wav'
  
  with speech_recognition.AudioFile(filename) as source:
    audio_new = r.record(source)
  
  words = r.recognize_google(audio_new, language="pl-PL")
  return words

def get_tarfile_name(checkpoint_folder):
    """Converts a folder path into a filename for a .tar archive"""
    tarfile_name = checkpoint_folder.replace(os.path.sep, '_') + '.tar'

    return tarfile_name

def unpack_checkpoints(run_name='run1'):
    """Copies the checkpoint folder from a mounted Google Drive."""
    
    checkpoint_folder = os.path.join('checkpoint', run_name)

    file_path = get_tarfile_name(checkpoint_folder)

    #shutil.copyfile(file_path, file_path)

    with tarfile.open(file_path, 'r') as tar:
        def is_within_directory(directory, target):
            
            abs_directory = os.path.abspath(directory)
            abs_target = os.path.abspath(target)
        
            prefix = os.path.commonprefix([abs_directory, abs_target])
            
            return prefix == abs_directory
        
        def safe_extract(tar, path=".", members=None, *, numeric_owner=False):
        
            for member in tar.getmembers():
                member_path = os.path.join(path, member.name)
                if not is_within_directory(path, member_path):
                    raise Exception("Attempted Path Traversal in Tar File")
        
            tar.extractall(path, members, numeric_owner=numeric_owner) 
            
        
        safe_extract(tar)
        
def cleaning_poem(poem):
    """aditional cleaning poems"""
    poem = poem.replace('— —','—')
    poem = poem.replace('  ',' ')
    poem = poem.replace('....','...')
    
    if len(poem)>1 and poem[-1] not in ['.','!','?']: 
        poem = poem[:-1] + '...'
    return poem

def averageLen(lst):
    """average list length"""
    lengths = [len(i) for i in lst]
    return 0 if len(lengths) == 0 else (float(sum(lengths)) / len(lengths)) 

def poems_check(poems, words):
    """for gpt return poems with checks"""
    poems_check = pd.DataFrame(columns=['lp', 'poem', 'check01', 'check02'])

    poems_list = list(map(str.strip, poems)) #usuwamy tzw. whitespaces
    words_list = words.split('\n') #rozdzielamy na osobne linie jeśli są

    for i in range(len(poems_list)):
        check01 = ':('
        check02 = ':('

        poem = poems_list[i].split('\n') #rozdzielamy pojedyńczy wiersz na linie
        if len(poem)>1: poem = poem[:-1] #usuwamy ostatni wers. często jest niepełny i ucięty

        #check01 - czy zaczyna się tak samo
        if fuzz.ratio(words_list[0], poem[0]) > 90: #words_list[0] == poem[0]:
            check01 = 'OK'

        #check02 - długość ostatniego wersa > śrenio połowa wersa
        if len(poem[-1]) > averageLen(poem)/2: 
            check02 = 'OK'
            
        poem = list(map(str.strip, poem))#usuwamy tzw. whitespaces
        poems_check.loc[i] = [i, cleaning_poem('\n'.join(poem)), check01, check02]

    return poems_check
    
def final_poem(poems, words):
    """return one poem"""
    final_poems_check = poems_check(poems, words)
    try:
      wygenerowany_wiersz = final_poems_check[(final_poems_check['check01']=='OK') & (final_poems_check['check02']=='OK')].sample(1).reset_index()['poem'][0]
    except:
      wygenerowany_wiersz = final_poems_check['poem'][0] #jak się nie uda weź pierwszy i już  

    return wygenerowany_wiersz
    
def ktora_epoka(wygenerowany_wiersz):

  dane = pd.DataFrame(columns=['wiersz'])
  dane.loc[0] = wygenerowany_wiersz 

  dane['wiersz_parsed_1'] = dane['wiersz'].str.replace("\r", " ")
  dane['wiersz_parsed_1'] = dane['wiersz_parsed_1'].str.replace("\n", " ")
  dane['wiersz_parsed_1'] = dane['wiersz_parsed_1'].str.replace("    ", " ")
  dane['wiersz_parsed_1'] = dane['wiersz_parsed_1'].str.lower()

  punctuation_signs = list("?:!.,;-")

  for punct_sign in punctuation_signs:
      dane['wiersz_parsed_1'] = dane['wiersz_parsed_1'].str.replace(punct_sign, '')

  nltk.download('stopwords')
  stop_words = get_stop_words('pl')

  for stop_word in stop_words:

      regex_stopword = r"\b" + stop_word + r"\b"
      dane['wiersz_parsed_1'] = dane['wiersz_parsed_1'].str.replace(regex_stopword, '')


  nlp = pl_spacy_model.load()

  nrows = len(dane)
  lemmatized_text_list = []

  for row in range(0, nrows):    
      print(row)
      lemmatized_list = []
      
      # Save the text and its words into an object
      text = dane.loc[row]['wiersz_parsed_1']

      # Iterate through every word to lemmatize
      doc = nlp(str(text))
      for token in doc:
          lemmatized_list.append(token.lemma_)
          
      # Join the list
      lemmatized_text = " ".join(lemmatized_list)
      
      # Append to the list containing the texts
      lemmatized_text_list.append(lemmatized_text)
      
  dane['wiersz_parsed_2'] = lemmatized_text_list

  tgidf_vocabulary = pd.read_pickle(r'tgidf.pickle')

  tgidf = TfidfTransformer()
  loaded_vec = CountVectorizer(decode_error="replace",vocabulary=tgidf_vocabulary)
  features_train = tgidf.fit_transform(loaded_vec.fit_transform(dane['wiersz_parsed_2']))

  model_wiersze = pd.read_pickle(r'best_xgb.pickle')

  y_pred = model_wiersze.predict_proba(features_train)

  category_codes = ['Barok', 'Młoda Polska', 'Oświecenie', 'Pozytywizm', 
                    'Renesans', 'Romantyzm', 'XX-lecie']

  wynik = pd.DataFrame(columns=['Epoka', 'Prawdopodobieństwo'])
  wynik['Epoka'] = category_codes
  wynik['Prawdopodobieństwo'] = np.concatenate(y_pred)

  plt.figure(figsize=(10,5))
 
  plt.bar(wynik['Epoka'], wynik['Prawdopodobieństwo']*100, color = '#969696')
    
  plt.xlabel('Epoka', fontsize=12, color='#323232')
  plt.ylabel('Prawdopodobieństwo', fontsize=12, color='#323232')
  plt.title('Z której epoki jest "Twój" wiersz?', fontsize=16, color='#323232')
  plt.gca().yaxis.set_major_formatter(mtick.PercentFormatter())
  
  plt.savefig('graph.jpg')
  plt.show()

  return wynik, dane

def df_to_table_new(slide, df, left, top, width, height, colnames=None):
  
    rows, cols = df.shape
    res = slide.shapes.add_table(rows + 1, cols, left, top, width, height)
    #res = slide.shapes.add_table(rows, cols, left, top, width, height)

    if colnames is None:
        colnames = list(df.columns)

    # Insert the column names
    for col_index, col_name in enumerate(colnames):
        # Column names can be tuples
        if not isinstance(col_name, str):
            col_name = " ".join(col_name)
        res.table.cell(0, col_index).text = col_name
        cell_new = res.table.cell(0, col_index)
        fill = cell_new.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(255, 255, 255)
        cell_new.text_frame.paragraphs[0].font.size = Pt(7.5)

    m = df.values

    for row in range(rows):
        for col in range(cols):
            val = m[row, col]
            text = str(val)
            res.table.cell(row + 1, col).text = text
            #res.table.cell(row, col).text = text
            cell_new = res.table.cell(row + 1, col)
            #cell_new = res.table.cell(row, col)
            fill = cell_new.fill
            fill.solid()
            fill.fore_color.rgb = RGBColor(255, 255, 255)
            run = res.table.cell(row + 1, col).text_frame.paragraphs[0].runs[0]
            run.font.italic = True
            run.font.size = Pt(16)

def pobierz_ppt(words, df_wiersz, ppt_name = 'wiersz_rezultat.pptx', img_name = 'graph.jpg'):
  genres = df_wiersz['wiersz'].str.split('\n')
  df = pd.DataFrame({
      'genres' : list(chain.from_iterable(genres.tolist()))
  })

  root = Presentation() 

  #Pierwszy slajd
  first_slide_layout = root.slide_layouts[5]  
  slide = root.slides.add_slide(first_slide_layout) 
  title = slide.shapes.title
  title.text = "\"" + words + "\""
  title.text_frame.paragraphs[0].font.color.rgb = RGBColor(59,89, 152)
  title.text_frame.paragraphs[0].alignment = PP_ALIGN.LEFT
  title.text_frame.paragraphs[0].font.size = Pt(25)
  title.text_frame.paragraphs[0].font.italic = True

  top = Inches(1.25)
  left = Inches(1.25)
  width = Inches(8.0)
  height = Inches(5.0)

  df.columns = ['']
  df_to_table_new(slide, df,left, top, width, height)

  #Drugi slajd
  third_slide_layout = root.slide_layouts[5]  
  slide = root.slides.add_slide(third_slide_layout) 
  title = slide.shapes.title
  title.text = "Twój wiersz przypomina najbardziej:"
  title.text_frame.paragraphs[0].font.color.rgb = RGBColor(59,89, 152)
  title.text_frame.paragraphs[0].alignment = PP_ALIGN.LEFT
  title.text_frame.paragraphs[0].font.size = Pt(25)
  title.text_frame.paragraphs[0].font.bold = True

  top = Inches(2.0)
  left = Inches(1.0)
  width = Inches(8.0)
  height = Inches(4.0)

  pic = slide.shapes.add_picture(img_name, left, top, width, height)

  #Zapisywanie
  root.save(ppt_name) 
  files.download(ppt_name) 

  print("done")